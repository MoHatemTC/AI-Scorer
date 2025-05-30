import requests
from fastapi import HTTPException
from typing import Union

from PyPDF2 import PdfReader
from openpyxl import load_workbook
from io import BytesIO
from docx import Document
from pptx import Presentation


def download_and_parse_file(url: str) -> Union[str, list[str]]:
    try:
        response = requests.get(url)
        response.raise_for_status()

        file_bytes = response.content
        file_text = response.text
        ext = url.lower().split("?")[0].split("#")[0].split(".")[-1]

        if ext in (
            "py",
            "js",
            "txt",
            "cpp",
            "c",
            "java",
        ):
            response.encoding = response.apparent_encoding
            return file_text.strip()

        # === Binary formats ===
        elif ext == "pdf":
            return parse_pdf(file_bytes)
        elif ext == "docx":
            return parse_docx(file_bytes)
        elif ext == "xlsx":
            return parse_xlsx(file_bytes)
        elif ext == "pptx":
            return parse_pptx(file_bytes)
        else:
            raise ValueError(f"Unsupported or unknown file extension: .{ext}")

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Download/Parsing error: {str(e)}")


def parse_pdf(file_bytes: bytes) -> str:
    pdf = PdfReader(BytesIO(file_bytes))
    text = ""
    for page in pdf.pages:
        text += page.extract_text() or ""
    return text.strip()


def parse_xlsx(file_bytes: bytes) -> list[str]:
    workbook = load_workbook(filename=BytesIO(file_bytes), data_only=True)
    results = []
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows(values_only=True):
            results.append(
                "\t".join([str(cell) if cell is not None else "" for cell in row])
            )
    return results  


def parse_docx(file_bytes: bytes) -> str:
    document = Document(BytesIO(file_bytes))
    return "\n".join([para.text for para in document.paragraphs if para.text.strip()])


def parse_pptx(file_bytes: bytes) -> str:
    presentation = Presentation(BytesIO(file_bytes))
    slides_text = []
    for slide in presentation.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        slides_text.append("\n".join(slide_text))
    return "\n\n".join(slides_text)
